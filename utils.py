import os
import cv2
import numpy as np
from moviepy.editor import ImageClip, VideoFileClip, AudioFileClip, CompositeVideoClip, concatenate_videoclips
from PIL import Image, ImageDraw, ImageFont
import tempfile
from pptx import Presentation
from gtts import gTTS
import azure.cognitiveservices.speech as speechsdk
from avatar_generator import create_avatar_video
import io

def extract_text_from_slide(slide):
    """Extract text content from a PowerPoint slide with better formatting"""
    text_content = []
    
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            # Check if it's a title
            if hasattr(shape, "is_title") and shape.is_title:
                text_content.append(f"Title: {shape.text.strip()}")
            else:
                # Handle bullet points
                if shape.text.strip():
                    points = shape.text.split('\n')
                    for point in points:
                        if point.strip():
                            text_content.append(point.strip())
    
    return " ".join(text_content)

def extract_text_from_slides(ppt_path):
    """Extract text from all slides in the presentation"""
    prs = Presentation(ppt_path)
    texts = []
    for slide in prs.slides:
        texts.append(extract_text_from_slide(slide))
    return texts

def convert_slide_to_image(ppt_path, slide_index, output_path, size=(1920, 1080)):
    """Convert a PowerPoint slide to an image using python-pptx and PIL"""
    try:
        # Ensure output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # Load the presentation
        prs = Presentation(ppt_path)
        
        # Get the slide
        slide = prs.slides[slide_index]
        
        # Create a blank image with white background
        img = Image.new('RGB', size, color='white')
        draw = ImageDraw.Draw(img)
        
        # Process each shape in the slide
        for shape in slide.shapes:
            # Get shape position and size
            left = int(shape.left * size[0] / 9144000)  # Convert EMU to pixels
            top = int(shape.top * size[1] / 6858000)
            width = int(shape.width * size[0] / 9144000)
            height = int(shape.height * size[1] / 6858000)
            
            # Handle different shape types
            if shape.shape_type == 13:  # Picture
                try:
                    image_stream = io.BytesIO(shape.image.blob)
                    pil_image = Image.open(image_stream)
                    pil_image = pil_image.resize((width, height), Image.Resampling.LANCZOS)
                    img.paste(pil_image, (left, top))
                except Exception as e:
                    print(f"Error processing image: {str(e)}")
            
            elif hasattr(shape, "text") and shape.text.strip():
                # Handle text with proper formatting
                text = shape.text.strip()
                
                # Determine if it's a title
                is_title = hasattr(shape, "is_title") and shape.is_title
                
                # Set font size based on whether it's a title and shape size
                if is_title:
                    font_size = int(min(width, height) / 15)
                else:
                    font_size = int(min(width, height) / 25)
                
                try:
                    font = ImageFont.truetype("arial.ttf", font_size)
                except:
                    font = ImageFont.load_default()
                
                # Draw text with word wrap
                words = text.split()
                lines = []
                current_line = []
                
                for word in words:
                    current_line.append(word)
                    line_text = ' '.join(current_line)
                    text_width = draw.textlength(line_text, font=font)
                    
                    if text_width > width - 20:  # Leave some padding
                        current_line.pop()
                        lines.append(' '.join(current_line))
                        current_line = [word]
                
                if current_line:
                    lines.append(' '.join(current_line))
                
                # Draw each line with proper spacing
                y = top
                line_height = font_size + 5
                for line in lines:
                    # Center text if it's a title
                    if is_title:
                        text_width = draw.textlength(line, font=font)
                        x = left + (width - text_width) / 2
                    else:
                        x = left + 10
                    
                    draw.text((x, y), line, font=font, fill='black')
                    y += line_height
            
            elif shape.shape_type == 1:  # Rectangle or other shape
                try:
                    # Draw shape outline
                    draw.rectangle([left, top, left + width, top + height], 
                                 outline='black', width=2)
                except Exception as e:
                    print(f"Error drawing shape: {str(e)}")
        
        # Save the image with high quality
        img.save(output_path, quality=95)
        return output_path
        
    except Exception as e:
        print(f"Error converting slide {slide_index}: {str(e)}")
        # Create a blank slide with error message
        img = Image.new('RGB', size, color='white')
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("arial.ttf", 40)
        except:
            font = ImageFont.load_default()
        draw.text((size[0]/2, size[1]/2), f"Slide {slide_index + 1}", font=font, fill='black', anchor="mm")
        img.save(output_path)
        return output_path

def generate_audio_from_text(text, output_path, voice_type="female", speech_rate=1.0, speech_pitch=0):
    """Generate audio from text using gTTS"""
    try:
        # Use gTTS for text-to-speech
        lang = "en"
        tld = "com" if voice_type.lower() == "female" else "co.uk"
        tts = gTTS(text=text, lang=lang, tld=tld, slow=(speech_rate < 1.0))
        tts.save(output_path)
        return output_path
            
    except Exception as e:
        print(f"Error in text-to-speech: {str(e)}")
        # Create a blank audio file if TTS fails
        from moviepy.editor import AudioClip
        
        # Create a silent audio clip of 1 second
        silent_clip = AudioClip(lambda t: 0, duration=1)
        silent_clip.write_audiofile(output_path)
        return output_path

def create_final_video(image_paths, audio_paths, avatar_paths, output_path):
    """Combine slides, audio, and avatars into final video"""
    clips = []
    for img_path, audio_path, avatar_path in zip(image_paths, audio_paths, avatar_paths):
        # Load components
        slide_clip = ImageClip(img_path)
        audio_clip = AudioFileClip(audio_path)
        avatar_clip = VideoFileClip(avatar_path)
        
        # Set durations
        slide_clip = slide_clip.set_duration(audio_clip.duration)
        avatar_clip = avatar_clip.set_duration(audio_clip.duration)
        
        # Position avatar in bottom right
        avatar_clip = avatar_clip.resize(width=slide_clip.w // 4)
        avatar_clip = avatar_clip.set_position((slide_clip.w - avatar_clip.w - 50, 50))
        
        # Combine clips
        final_clip = CompositeVideoClip([slide_clip, avatar_clip])
        final_clip = final_clip.set_audio(audio_clip)
        clips.append(final_clip)
    
    # Concatenate all clips
    final_video = concatenate_videoclips(clips)
    final_video.write_videofile(output_path, fps=30)
    return output_path

def process_presentation(ppt_path, output_path, voice_type="female", language="en", tld="com", speech_rate=1.0, speech_pitch=0):
    """
    Process the entire presentation with enhanced features
    """
    try:
        # Extract text from slides
        slide_texts = extract_text_from_slides(ppt_path)
        
        # Convert slides to images
        image_paths = []
        for i in range(len(slide_texts)):
            image_path = convert_slide_to_image(ppt_path, i)
            if image_path:
                image_paths.append(image_path)
        
        # Generate audio for each slide
        audio_paths = []
        for text in slide_texts:
            audio_path = generate_audio_from_text(text, voice_type, language, tld, speech_rate, speech_pitch)
            if audio_path:
                audio_paths.append(audio_path)
        
        # Create avatar videos for each slide
        avatar_paths = []
        for text in slide_texts:
            avatar_path = create_avatar_video(text, voice_type)
            if avatar_path:
                avatar_paths.append(avatar_path)
        
        # Create final video
        create_final_video(image_paths, audio_paths, avatar_paths, output_path)
        
        return output_path
        
    except Exception as e:
        print(f"Error processing presentation: {str(e)}")
        return None 